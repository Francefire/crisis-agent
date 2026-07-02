"""
Build a 5-slide Day-1 deck (Ultia x CNC) as a .pptx that imports cleanly into Canva.
Run: .venv/bin/python analysis/build_slides.py
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.dirname(HERE)
S = json.load(open(os.path.join(HERE, "out", "key_figures.json")))

PURPLE = RGBColor(0x5B, 0x2B, 0xE0)
DARK   = RGBColor(0x2A, 0x10, 0x5C)
TEAL   = RGBColor(0x2B, 0xBF, 0xA0)
ORANGE = RGBColor(0xE0, 0x66, 0x2B)
GREY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF3, 0xF0, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Arial"
    return tb


def kicker(slide, txt, color=PURPLE):
    text(slide, Inches(0.6), Inches(0.35), Inches(9), Inches(0.4),
         [(txt.upper(), 13, color, True)])


def title(slide, txt, color=DARK):
    text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(1.0),
         [(txt, 30, color, True)])


def footer(slide, n):
    text(slide, Inches(11.6), Inches(7.0), Inches(1.6), Inches(0.4),
         [(f"NEXA · Datathon   {n}/5", 9, GREY, False)], align=PP_ALIGN.RIGHT)


def pic_fit(slide, path, x, y, w, h):
    """Add picture fit into (w,h) keeping aspect, centered."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    tw, th = w, int(w / ar)
    if th > h:
        th = h; tw = int(h * ar)
    px = x + (w - tw) // 2
    py = y + (h - th) // 2
    slide.shapes.add_picture(path, px, py, tw, th)


# ---- helpers on data
et = S["engagement_type_counts"]
rt_pct = S["retweet_share_pct"]
sent = S["sentiment_counts"]
neg_pct = round(sent["negative"] / S["n_rows"] * 100)
peak = S["peak_day"]; peakv = S["peak_day_volume"]
amp = list(S["top15_amplified_authors"].items())
terms = list(S["top40_terms_docfreq"].items())

# ============================================================ SLIDE 1
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.5),
     [("DATATHON · 2 JOURS", 14, TEAL, True)])
box(s, Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.7), fill=PURPLE)
text(s, Inches(1.4), Inches(2.4), Inches(10.5), Inches(1.7),
     [("Affaire Ultia × CNC : anatomie d'un déferlement viral", 32, WHITE, True)],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(1.2), Inches(4.5), Inches(10.9), Inches(1.6),
     [("Un fait réel, transformé en tempête sur X en 4 jours.", 18, WHITE, False),
      ("Fait  →  Narratif  →  Dynamique : nous analysons la dynamique.",
       16, TEAL, True)], align=PP_ALIGN.CENTER, space_after=10)
text(s, Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.6),
     [(f"{S['n_rows']:,} messages · {S['n_authors']:,} comptes · "
       f"{S['date_min'][:10]} → {S['date_max'][:10]}".replace(",", " "),
       15, WHITE, True)], align=PP_ALIGN.CENTER)

# ============================================================ SLIDE 2
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
kicker(s, "Propagation — Timeline")
title(s, "Une déflagration de 4 jours, portée par le retweet")
pic_fit(s, os.path.join(FIG, "01_timeline_daily.png"),
        Inches(0.5), Inches(1.7), Inches(8.4), Inches(4.9))
# right stat cards
cards = [
    (f"{rt_pct:.0f}%", "des messages sont des retweets\n(amplification, pas débat)", PURPLE),
    (f"{peakv:,}".replace(",", " "), f"messages le {peak}\n(jour de pic)", ORANGE),
    (f"{et['ORIGINAL']}", "messages originaux seulement\n(1,9 % du corpus)", TEAL),
]
cy = Inches(1.8)
for big, small, col in cards:
    box(s, Inches(9.2), cy, Inches(3.6), Inches(1.5), fill=LIGHT)
    text(s, Inches(9.4), cy + Inches(0.12), Inches(3.2), Inches(0.7),
         [(big, 30, col, True)])
    text(s, Inches(9.4), cy + Inches(0.78), Inches(3.2), Inches(0.7),
         [(small, 12, GREY, False)])
    cy += Inches(1.65)
footer(s, 2)

# ============================================================ SLIDE 3
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
kicker(s, "Acteurs")
title(s, "Deux couches : des sources amplifiées, une piétaille de relais")
pic_fit(s, os.path.join(FIG, "03_acteurs_top.png"),
        Inches(0.4), Inches(1.9), Inches(8.6), Inches(4.7))
box(s, Inches(9.2), Inches(1.9), Inches(3.7), Inches(4.6), fill=LIGHT)
lines = [("Top comptes amplifiés (patient zéro / relais)", 13, DARK, True)]
for a, c in amp[:8]:
    lines.append((f"@{a}", 12, PURPLE, True))
    lines.append((f"   {c:,} reprises".replace(",", " "), 11, GREY, False))
text(s, Inches(9.4), Inches(2.05), Inches(3.35), Inches(4.3), lines, space_after=2)
footer(s, 3)

# ============================================================ SLIDE 4
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
kicker(s, "Narratifs & Sémantique")
title(s, "Trois cadrages qui s'affrontent — un ton majoritairement négatif")
pic_fit(s, os.path.join(FIG, "07_narratifs_top_terms.png"),
        Inches(0.3), Inches(1.7), Inches(5.4), Inches(5.4))
frames = [
    ("L'argent public", "cnc · fonds · subventions · millions · euros · aides", ORANGE),
    ("Les camps politiques", "gauche · extrême · droite · français", PURPLE),
    ("Qui finance-t-on ?", "talent · streameuse · créateurs · cinéma · film", TEAL),
]
fy = Inches(1.9)
for t, w, col in frames:
    box(s, Inches(6.0), fy, Inches(6.9), Inches(1.15), fill=LIGHT)
    text(s, Inches(6.2), fy + Inches(0.1), Inches(6.5), Inches(1.0),
         [(t, 16, col, True), (w, 12, GREY, False)], space_after=4)
    fy += Inches(1.3)
box(s, Inches(6.0), fy, Inches(6.9), Inches(1.2), fill=DARK)
text(s, Inches(6.2), fy + Inches(0.12), Inches(6.5), Inches(1.0),
     [(f"Sentiment : {neg_pct}% négatif · {round(sent['neutral']/S['n_rows']*100)}% neutre "
       f"· {round(sent['positive']/S['n_rows']*100)}% positif",
       15, WHITE, True),
      ("La négativité culmine avec le pic de volume.", 12, TEAL, False)],
     space_after=4)
footer(s, 4)

# ============================================================ SLIDE 5
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
kicker(s, "Coordination & suite", TEAL)
text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(1.0),
     [("Signaux de coordination → un agent IA pour outiller la riposte", 28, WHITE, True)])
# coordination signals
box(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.6), fill=PURPLE)
maxdup = max(S["top20_duplicated_messages"].values())
text(s, Inches(0.85), Inches(2.2), Inches(5.5), Inches(4.3),
     [("Ce que montre la coordination", 17, WHITE, True),
      (f"• Messages copiés-collés jusqu'à {maxdup:,} fois à l'identique".replace(",", " "), 14, WHITE, False),
      ("• Rafales synchronisées à la même minute", 14, WHITE, False),
      (f"• {S['n_authors_1_msg']:,} comptes ne postent qu'1 message ".replace(",", " ")
       + "(masse de relais)", 14, WHITE, False),
      (f"• {S['verified_share_pct']}% de comptes certifiés seulement", 14, WHITE, False)],
     space_after=10)
# next step / agent
box(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.6), fill=TEAL)
text(s, Inches(7.15), Inches(2.2), Inches(5.3), Inches(4.3),
     [("Prochaine étape : l'agent « Analyste »", 17, DARK, True),
      ("Modèle + outil d'accès au corpus + mémoire", 14, DARK, True),
      ("• Interroge le corpus (volume, comptes, narratifs)", 13, DARK, False),
      ("• Résume le narratif dominant d'un échantillon", 13, DARK, False),
      ("• Détecte les signaux de coordination", 13, DARK, False),
      ("Réutilisable pour d'autres crises informationnelles.", 13, DARK, True)],
     space_after=8)
footer(s, 5)

path = os.path.join(OUT, "Datathon_Ultia_CNC_Jour1.pptx")
prs.save(path)
print("Saved", path)
